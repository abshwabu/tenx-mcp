from pptx import Presentation
from pptx.util import Inches, Pt
import os

def create_presentation():
    prs = Presentation()

    # Slide 1: Title
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "MCP Setup Challenge"
    subtitle.text = "Configuring AI Agent Workflows and MCP Integration\nPrepared by Gemini CLI Agent"

    # Slide 2: Research & Methodology
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Research & Methodology"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Workflow Analysis"
    p = tf.add_paragraph()
    p.text = "Analyzed Boris Cherny's (Claude Code) workflow principles."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Focus on 'Plan vs. Act' modes and rigorous codebase exploration."
    p.level = 1

    # Slide 3: Configuration Implementation
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Configuration Implementation"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = ".vscode/mcp.json"
    p = tf.add_paragraph()
    p.text = "Configured proxy URL: https://mcppulse.10academy.org/proxy"
    p.level = 1
    p = tf.add_paragraph()
    p.text = ".github/copilot-instructions.md"
    p.level = 0
    p.text = "Structured rules for AI behavior, safety, and mandatory logging."

    # Slide 4: Proof of Concept (PoC)
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "MCP Proof of Concept"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Demonstrated Functionality"
    p = tf.add_paragraph()
    p.text = "Implemented local MCP server using FastMCP."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Validated tool listing and execution via stdio client."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Confirmed protocol compliance (JSON-RPC over stdio)."
    p.level = 1

    # Slide 5: Conclusion
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Next Steps"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Future Enhancements"
    p = tf.add_paragraph()
    p.text = "Complete OAuth2 authentication flow for remote proxy."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Expand toolset within the Tenx MCP server."
    p.level = 1

    output_path = 'deliverables/MCP_Setup_Challenge.pptx'
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    create_presentation()
